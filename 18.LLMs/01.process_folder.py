import os
import json
import pdfplumber
import pytesseract
from PIL import Image
import pandas as pd
from datetime import datetime
from docx import Document
from pptx import Presentation
import warnings
import re

# Suppress specific CropBox warning from pdfplumber
import logging
logging.getLogger("pdfminer").setLevel(logging.ERROR)

class CropBoxFilter:
    def filter(self, record):
        return not (record.levelno == logging.WARNING and 'CropBox missing from /Page, defaulting to MediaBox' in record.getMessage())

logging.getLogger().addFilter(CropBoxFilter())

# Add tkinter for folder selection
def select_folder():
    import tkinter as tk
    from tkinter import filedialog
    root = tk.Tk()
    root.withdraw()
    folder_selected = filedialog.askdirectory(title='Select folder to process')
    root.destroy()
    return folder_selected

def extract_text_from_pdf(pdf_path):
    text = ""
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    except Exception:
        pass
    return text.strip()

def ocr_pdf(pdf_path):
    text = ""
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                img = page.to_image(resolution=300).original
                pil_img = Image.fromarray(img)
                page_text = pytesseract.image_to_string(pil_img)
                if page_text:
                    text += page_text + "\n"
    except Exception:
        pass
    return text.strip()

def extract_text_from_excel(excel_path):
    text = ""
    try:
        xls = pd.ExcelFile(excel_path)
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name, dtype=str)
            text += f"\n--- Sheet: {sheet_name} ---\n"
            text += df.fillna('').to_string(index=False, header=True)
    except Exception:
        pass
    return text.strip()

def extract_text_from_docx(docx_path):
    text = ""
    try:
        doc = Document(docx_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception:
        pass
    return text.strip()

def extract_text_from_pptx(pptx_path):
    text = ""
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception:
        pass
    return text.strip()

def process_folder(folder_path, output_json_path):
    data = []
    for root, dirs, files in os.walk(folder_path):
        for file in files:
            print(f"[INFO] Processing file: {file}")
            file_path = os.path.join(root, file)
            ext = file.lower().split('.')[-1]
            entry = {
                'file_name': file,
                'file_path': file_path,
                'type': ext,
                'extracted_at': datetime.now().isoformat(),
                'text': ''
            }
            reason = None
            if ext == 'pdf':
                text = extract_text_from_pdf(file_path)
                if not text:
                    text = ocr_pdf(file_path)
                    if text:
                        reason = 'extracted via OCR'
                    else:
                        reason = 'no text extracted (PDF and OCR failed)'
                else:
                    reason = 'extracted as text PDF'
                entry['text'] = text
            elif ext in ['xls', 'xlsx']:
                text = extract_text_from_excel(file_path)
                entry['text'] = text
                reason = 'extracted from Excel' if text else 'no text extracted (Excel)'
            elif ext == 'docx':
                text = extract_text_from_docx(file_path)
                entry['text'] = text
                reason = 'extracted from Word' if text else 'no text extracted (Word)'
            elif ext == 'pptx':
                text = extract_text_from_pptx(file_path)
                entry['text'] = text
                reason = 'extracted from PowerPoint' if text else 'no text extracted (PowerPoint)'
            else:
                print(f"[WARN] Skipped unsupported file type: {file}")
                continue
            if entry['text']:
                data.append(entry)
            else:
                print(f"[WARN] No text extracted from: {file} ({reason})")
    with open(output_json_path, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    print(f"Processed {len(data)} documents. Output saved to {output_json_path}")

if __name__ == "__main__":
    # New output directory outside the repo
    output_dir = r'C:/Users/Omar Essam2/OneDrive - Rowad Modern Engineering/x004 Data Science/03.rme.db/05.llm/extracted_json'
    os.makedirs(output_dir, exist_ok=True)
    folder = select_folder()
    if not folder:
        print("No folder selected. Exiting.")
    else:
        folder_name = os.path.basename(os.path.normpath(folder))
        output_json = os.path.join(output_dir, f"{folder_name}_extracted.json")
        process_folder(folder, output_json) 
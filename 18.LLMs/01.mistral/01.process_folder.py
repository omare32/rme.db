import os
import json
import pdfplumber
import pytesseract
from PIL import Image
import pandas as pd
from datetime import datetime
from docx import Document
from pptx import Presentation
import warnings
import re
import zipfile
import rarfile
import tempfile
import shutil

# Suppress specific CropBox warning from pdfplumber
import logging
logging.getLogger("pdfminer").setLevel(logging.ERROR)

class CropBoxFilter:
    def filter(self, record):
        return not (record.levelno == logging.WARNING and 'CropBox missing from /Page, defaulting to MediaBox' in record.getMessage())

logging.getLogger().addFilter(CropBoxFilter())

# Add tkinter for folder selection
def select_folder():
    import tkinter as tk
    from tkinter import filedialog
    root = tk.Tk()
    root.withdraw()
    folder_selected = filedialog.askdirectory(title='Select folder to process')
    root.destroy()
    return folder_selected

def extract_text_from_pdf(pdf_path):
    text = ""
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    except Exception:
        pass
    return text.strip()

def ocr_pdf(pdf_path):
    text = ""
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                img = page.to_image(resolution=300).original
                pil_img = Image.fromarray(img)
                page_text = pytesseract.image_to_string(pil_img)
                if page_text:
                    text += page_text + "\n"
    except Exception:
        pass
    return text.strip()

def extract_project_name(file_path):
    # Try to extract project name from folder path or filename
    path_parts = os.path.normpath(file_path).split(os.sep)
    # Remove common non-project folders
    ignore = set(['', '2024', '2025', '2023', '2022', '2021', '2020', 'Mar. 2024', 'Jan-25', '3-Mar-2025', '01 Jan-25'])
    # Look for a likely project name in the path
    for part in reversed(path_parts[:-1]):
        if part not in ignore and not re.match(r'\d{4}', part):
            return part
    # Try to extract from filename
    fname = os.path.basename(file_path)
    match = re.search(r'([A-Za-z0-9\- ]+Project|[A-Za-z0-9\- ]+Dashboard|[A-Za-z0-9\- ]+Lot|[A-Za-z0-9\- ]+Port|[A-Za-z0-9\- ]+Sector)', fname)
    if match:
        return match.group(0)
    # If nothing found, label as GLOBAL
    return 'GLOBAL'

def extract_text_from_excel(excel_path):
    text = ""
    try:
        ext = excel_path.lower().split('.')[-1]
        if ext == 'xlsb':
            xls = pd.ExcelFile(excel_path, engine='pyxlsb')
        else:
            xls = pd.ExcelFile(excel_path)
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name, dtype=str)
            # Only extract if <1000 rows and <30 columns and has headers
            if df.shape[0] < 1000 and df.shape[1] < 30 and all(df.columns.str.strip() != ''):
                text += f"\n--- Sheet: {sheet_name} ---\n"
                text += df.fillna('').to_csv(index=False)
    except Exception:
        pass
    return text.strip()

def extract_text_from_docx(docx_path):
    text = ""
    try:
        doc = Document(docx_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception:
        pass
    return text.strip()

def extract_text_from_pptx(pptx_path):
    text = ""
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception:
        pass
    return text.strip()

def extract_text_from_ppt(file_path):
    """Attempt to convert .ppt to .pptx using unoconv, then extract text."""
    import subprocess
    import tempfile
    text = ""
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = os.path.join(tmpdir, os.path.basename(file_path) + '.pptx')
            # Try to convert .ppt to .pptx using unoconv
            result = subprocess.run([
                'unoconv', '-f', 'pptx', '-o', pptx_path, file_path
            ], capture_output=True)
            if result.returncode == 0 and os.path.exists(pptx_path):
                text = extract_text_from_pptx(pptx_path)
            else:
                print(f"[WARN] Failed to convert .ppt to .pptx: {file_path}")
    except Exception as e:
        print(f"[WARN] Exception during .ppt conversion: {file_path} | {e}")
    return text.strip()

def extract_text_from_pbix(pbix_path):
    """Try to extract text from .pbix by unzipping and reading .json/.xml files inside."""
    import zipfile
    import tempfile
    text = ""
    try:
        with zipfile.ZipFile(pbix_path, 'r') as z:
            with tempfile.TemporaryDirectory() as tmpdir:
                z.extractall(tmpdir)
                for root, _, files in os.walk(tmpdir):
                    for file in files:
                        if file.endswith('.json') or file.endswith('.xml'):
                            try:
                                with open(os.path.join(root, file), 'r', encoding='utf-8', errors='ignore') as f:
                                    content = f.read()
                                    # Only keep text if it's not too large
                                    if len(content) < 500_000:
                                        text += f"\n--- {file} ---\n" + content
                            except Exception:
                                continue
    except Exception as e:
        print(f"[WARN] Failed to extract .pbix: {pbix_path} | {e}")
    return text.strip()

def extract_text_from_doc(doc_path):
    """Convert .doc to .docx using unoconv, then extract text."""
    import subprocess
    import tempfile
    import os
    text = ""
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            docx_path = os.path.join(tmpdir, os.path.basename(doc_path) + '.docx')
            result = subprocess.run([
                'unoconv', '-f', 'docx', '-o', docx_path, doc_path
            ], capture_output=True)
            if result.returncode == 0 and os.path.exists(docx_path):
                text = extract_text_from_docx(docx_path)
            else:
                print(f"[WARN] Failed to convert .doc to .docx: {doc_path}")
    except Exception as e:
        print(f"[WARN] Exception during .doc conversion: {doc_path} | {e}")
    return text.strip()

def process_file(file_path, data):
    ext = file_path.lower().split('.')[-1]
    project = extract_project_name(file_path)
    print(f"[INFO] Processing file: {os.path.basename(file_path)} | Project: {project}")
    entry = {
        'file_name': os.path.basename(file_path),
        'file_path': file_path,
        'type': ext,
        'project': project,
        'extracted_at': datetime.now().isoformat(),
        'text': ''
    }
    reason = None
    if ext == 'pdf':
        text = extract_text_from_pdf(file_path)
        if not text:
            text = ocr_pdf(file_path)
            if text:
                reason = 'extracted via OCR'
            else:
                reason = 'no text extracted (PDF and OCR failed)'
        else:
            reason = 'extracted as text PDF'
        entry['text'] = text
    elif ext in ['xls', 'xlsx', 'xlsb']:
        text = extract_text_from_excel(file_path)
        entry['text'] = text
        reason = 'extracted from Excel' if text else 'no text extracted (Excel)'
    elif ext == 'docx':
        text = extract_text_from_docx(file_path)
        entry['text'] = text
        reason = 'extracted from Word' if text else 'no text extracted (Word)'
    elif ext in ['pptx', 'pptm']:
        text = extract_text_from_pptx(file_path)
        entry['text'] = text
        reason = 'extracted from PowerPoint' if text else 'no text extracted (PowerPoint)'
    elif ext == 'ppt':
        text = extract_text_from_ppt(file_path)
        entry['text'] = text
        reason = 'extracted from PPT (via conversion)' if text else 'no text extracted (PPT)'
    elif ext == 'pbix':
        text = extract_text_from_pbix(file_path)
        entry['text'] = text
        reason = 'extracted from PBIX' if text else 'no text extracted (PBIX)'
    elif ext == 'doc':
        text = extract_text_from_doc(file_path)
        entry['text'] = text
        reason = 'extracted from DOC (via conversion)' if text else 'no text extracted (DOC)'
    else:
        print(f"[WARN] Skipped unsupported file type: {os.path.basename(file_path)}")
        return
    if entry['text']:
        data.append(entry)
    else:
        print(f"[WARN] No text extracted from: {os.path.basename(file_path)} ({reason})")

def process_folder(folder_path, output_json_path):
    data = []
    for root, dirs, files in os.walk(folder_path):
        for file in files:
            file_path = os.path.join(root, file)
            ext = file.lower().split('.')[-1]
            if ext in ['zip', 'rar']:
                print(f"[INFO] Extracting archive: {file}")
                with tempfile.TemporaryDirectory() as tmpdir:
                    try:
                        if ext == 'zip':
                            with zipfile.ZipFile(file_path, 'r') as z:
                                z.extractall(tmpdir)
                        elif ext == 'rar':
                            with rarfile.RarFile(file_path, 'r') as r:
                                r.extractall(tmpdir)
                        # Recursively process extracted files
                        for subroot, _, subfiles in os.walk(tmpdir):
                            for subfile in subfiles:
                                subfile_path = os.path.join(subroot, subfile)
                                process_file(subfile_path, data)
                    except Exception as e:
                        print(f"[WARN] Failed to extract {file}: {e}")
            else:
                process_file(file_path, data)
    with open(output_json_path, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    print(f"Processed {len(data)} documents. Output saved to {output_json_path}")

if __name__ == "__main__":
    # New output directory outside the repo
    output_dir = r'D:/OEssam/extracted_json'
    os.makedirs(output_dir, exist_ok=True)
    folder = select_folder()
    if not folder:
        print("No folder selected. Exiting.")
    else:
        folder_name = os.path.basename(os.path.normpath(folder))
        output_json = os.path.join(output_dir, f"{folder_name}_extracted.json")
        process_folder(folder, output_json) 
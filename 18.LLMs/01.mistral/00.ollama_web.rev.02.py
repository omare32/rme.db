from flask import Flask, render_template_string, request, redirect, url_for
import requests
import pdfplumber
import pytesseract
from PIL import Image
import io
import os
from docx import Document
from pptx import Presentation
import pandas as pd
from dotenv import load_dotenv
import openai

load_dotenv()
openai.api_key = os.getenv('OPENAI_API_KEY')

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = r'D:/OEssam/uploads'
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

HTML = '''
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>RME Chatbot Rev003a</title>
    <style>
        body { font-family: Arial, sans-serif; background: #f4f6fb; margin: 0; padding: 0; }
        .container { max-width: 700px; margin: 40px auto; background: #fff; border-radius: 10px; box-shadow: 0 2px 8px #0001; padding: 32px; }
        h1 { text-align: center; color: #333; }
        .model-info { text-align: center; background: #e0e7ff; color: #222; border-radius: 6px; padding: 8px 0; margin-bottom: 24px; font-weight: bold; }
        textarea { width: 100%; min-height: 80px; padding: 10px; border-radius: 6px; border: 1px solid #ccc; font-size: 1rem; }
        .btn-row { display: flex; gap: 12px; margin-top: 12px; }
        button { background: #4f8cff; color: #fff; border: none; padding: 12px 24px; border-radius: 6px; font-size: 1rem; cursor: pointer; }
        button:hover { background: #2563eb; }
        .mixtral-btn { background: #22c55e; }
        .mixtral-btn:hover { background: #15803d; }
        .gpt-btn { background: #ffd700; color: #222; font-weight: bold; }
        .gpt-btn:hover { background: #ffb300; }
        .response { margin-top: 24px; background: #f0f4fa; padding: 16px; border-radius: 6px; white-space: pre-wrap; font-size: 1.05rem; }
        label { font-weight: bold; }
        .pdf-section { margin-top: 32px; }
        .pdf-text { background: #f9fafb; border: 1px solid #ddd; border-radius: 6px; padding: 12px; margin-top: 8px; max-height: 200px; overflow-y: auto; font-size: 0.98rem; display: none; }
        .show { display: block !important; }
    </style>
    <script>
        function togglePdfText() {
            var pdfTextDiv = document.getElementById('pdfTextDiv');
            if (pdfTextDiv.classList.contains('show')) {
                pdfTextDiv.classList.remove('show');
            } else {
                pdfTextDiv.classList.add('show');
            }
        }
    </script>
</head>
<body>
    <div class="container">
        <h1>RME Chatbot Rev003a</h1>
        <div class="model-info">Choose a model to answer your question</div>
        <form method="post" enctype="multipart/form-data">
            <label for="file">Upload a file (PDF, Excel, Word, PowerPoint):</label><br>
            <input type="file" id="file" name="file" accept=".pdf,.xls,.xlsx,.docx,.pptx"><br><br>
            <label for="prompt">Enter your prompt:</label><br>
            <textarea id="prompt" name="prompt" required>{{ prompt or '' }}</textarea><br>
            <div class="btn-row">
                <button type="submit" name="model" value="mistral">Send to Mistral</button>
                <button type="submit" name="model" value="mixtral" class="mixtral-btn">Send to Mixtral</button>
                <button type="submit" name="model" value="gpt" class="gpt-btn">Send to GPT (💲)</button>
            </div>
        </form>
        {% if response %}
        <div class="response">
            <strong>Response:</strong><br>
            {{ response }}
        </div>
        {% endif %}
        {% if file_text %}
        <div class="pdf-section">
            <button type="button" onclick="togglePdfText()">Show Extracted File Text</button>
            <div id="pdfTextDiv" class="pdf-text">{{ file_text }}</div>
        </div>
        {% endif %}
    </div>
</body>
</html>
'''

def get_mistral_version():
    try:
        r = requests.get('http://localhost:11434/api/tags', timeout=5)
        if r.status_code == 200:
            models = r.json().get('models', [])
            for model in models:
                if model['name'].startswith('mistral'):
                    details = model.get('details', {})
                    param_size = details.get('parameter_size')
                    if param_size:
                        return param_size
                    if '7b' in model['name'].lower():
                        return '7B'
            return '7B'
        else:
            return '7B'
    except Exception:
        return '7B'

def extract_text_from_pdf(pdf_path):
    text = ""
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    except Exception:
        pass
    return text.strip()

def ocr_pdf(pdf_path):
    text = ""
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                img = page.to_image(resolution=300).original
                pil_img = Image.fromarray(img)
                page_text = pytesseract.image_to_string(pil_img)
                if page_text:
                    text += page_text + "\n"
    except Exception:
        pass
    return text.strip()

def extract_text_from_excel(excel_path):
    text = ""
    try:
        xls = pd.ExcelFile(excel_path)
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name, dtype=str)
            text += f"\n--- Sheet: {sheet_name} ---\n"
            text += df.fillna('').to_string(index=False, header=True)
    except Exception:
        pass
    return text.strip()

def extract_text_from_docx(docx_path):
    text = ""
    try:
        doc = Document(docx_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception:
        pass
    return text.strip()

def extract_text_from_pptx(pptx_path):
    text = ""
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception:
        pass
    return text.strip()

def ask_ollama_model(prompt, model_name):
    try:
        r = requests.post(
            'http://localhost:11434/api/generate',
            json={
                'model': model_name,
                'prompt': prompt,
                'stream': False
            },
            timeout=120
        )
        if r.status_code == 200:
            return r.json().get('response', '(No response)')
        else:
            return f"Ollama API error: {r.status_code}"
    except Exception as e:
        return f"Error: {e}"

def ask_gpt(prompt):
    try:
        resp = openai.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=512,
            temperature=0.2
        )
        return resp.choices[0].message.content.strip()
    except Exception as e:
        return f"OpenAI API error: {e}"

@app.route('/', methods=['GET', 'POST'])
def index():
    response = None
    prompt = None
    file_text = None
    if request.method == 'POST':
        prompt = request.form['prompt']
        uploaded_file = request.files.get('file')
        model_choice = request.form.get('model', 'mistral')
        if uploaded_file and uploaded_file.filename:
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], uploaded_file.filename)
            uploaded_file.save(file_path)
            ext = uploaded_file.filename.lower().split('.')[-1]
            if ext == 'pdf':
                file_text = extract_text_from_pdf(file_path)
                if not file_text:
                    file_text = ocr_pdf(file_path)
            elif ext in ['xls', 'xlsx']:
                file_text = extract_text_from_excel(file_path)
            elif ext == 'docx':
                file_text = extract_text_from_docx(file_path)
            elif ext == 'pptx':
                file_text = extract_text_from_pptx(file_path)
        if file_text:
            full_prompt = f"Context from file:\n{file_text}\n\nQuestion: {prompt}"
        else:
            full_prompt = prompt
        if model_choice == 'mistral':
            response = ask_ollama_model(full_prompt, 'mistral:latest')
        elif model_choice == 'mixtral':
            response = ask_ollama_model(full_prompt, 'dolphin-mixtral:latest')
        elif model_choice == 'gpt':
            response = ask_gpt(full_prompt)
        else:
            response = 'Unknown model selected.'
    return render_template_string(HTML, response=response, prompt=prompt, file_text=file_text)

if __name__ == '__main__':
    app.run(host="0.0.0.0", debug=True) 